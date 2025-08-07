import os
import pickle
import logging
from typing import Dict, List, Optional, Tuple, Any
from dotenv import load_dotenv

# LangChain imports
from langchain_community.document_loaders import (
    PyPDFLoader,
    TextLoader,
    Docx2txtLoader,
    CSVLoader,
    UnstructuredPowerPointLoader
)
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_groq import ChatGroq
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.output_parsers import StrOutputParser
from langchain_core.runnables import RunnablePassthrough
from langchain_core.documents import Document

# PDF and document processing
import fitz  # PyMuPDF
from pptx import Presentation
import re
import hashlib
import pathlib

# Import new components
from conversation_manager import ConversationManager
from router import QueryRouter
from web_search import WebSearchClient
from prompts import ANSWER_PROMPT, ANSWER_CORE_INSTRUCTIONS

# Load environment variables
load_dotenv()

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

class CourseRAG:
    """RAG system for university course materials using Groq API"""
    
    def __init__(self, course_id: str, materials_dir: Optional[str] = None,
                discipline: Optional[str] = None):
        """Initialize the RAG system for a specific course
        
        Args:
            course_id: Identifier for the course (e.g., "CS101")
            materials_dir: Directory containing course materials (optional)
            discipline: Course discipline for prompt templates (optional)
        """
        self.course_id = course_id
        self.materials_dir = materials_dir or os.getenv("MATERIALS_DIR", "course_materials")
        self.course_dir = os.path.join(self.materials_dir, course_id)
        self.vectorstore_dir = "vectorstores"
        self.vectorstore_path = os.path.join(self.vectorstore_dir, f"{course_id}_vectorstore.pkl")
        self.discipline = discipline
        
        # Initialize embeddings
        self.embeddings = self._create_embeddings()
        
        # Initialize vector store
        self.vectorstore = self._initialize_vectorstore()
        
        # Initialize LLM
        self.llm = self._initialize_llm()
        
        # Initialize conversation manager
        self.conversation_manager = ConversationManager()
        
        # Build the RAG chain
        self.rag_chain = self._build_rag_chain()

        # Agentic components (lazy)
        self._router: Optional[QueryRouter] = None
        self._web: Optional[WebSearchClient] = None

    @staticmethod
    def _make_json_serializable(obj: Any) -> Any:
        """Convert NumPy types and nested structures to Python native types."""
        try:
            import numpy as _np  # local import to avoid hard dependency elsewhere
        except Exception:  # pragma: no cover
            _np = None

        if _np is not None and isinstance(obj, (_np.float32, _np.float64)):
            return float(obj)
        if _np is not None and isinstance(obj, (_np.int32, _np.int64)):
            return int(obj)
        if _np is not None and isinstance(obj, _np.ndarray):
            return obj.tolist()
        if isinstance(obj, dict):
            return {k: CourseRAG._make_json_serializable(v) for k, v in obj.items()}
        if isinstance(obj, list):
            return [CourseRAG._make_json_serializable(v) for v in obj]
        return obj
    
    def _create_embeddings(self):
        """Create embeddings using a local HuggingFace model"""
        logger.info("Initializing embeddings model...")
        return HuggingFaceEmbeddings(
            model_name="all-MiniLM-L6-v2",
            model_kwargs={'device': 'cpu'}
        )
    
    def _initialize_llm(self):
        """Initialize the Groq LLM"""
        groq_api_key = os.getenv("GROQ_API_KEY")
        if not groq_api_key:
            raise ValueError("GROQ_API_KEY not found in environment variables")
        
        model_name = os.getenv("GROQ_MODEL", "llama3-70b-8192")
        logger.info(f"Initializing Groq LLM with model: {model_name}")
        
        return ChatGroq(
            api_key=groq_api_key,
            model_name=model_name
        )
    
    def _load_documents(self) -> List[Document]:
        """Load all documents from the course materials directory"""
        logger.info(f"Loading documents for course {self.course_id}")
        
        documents = []
        course_dir = os.path.join(self.materials_dir, self.course_id)
        
        if not os.path.exists(course_dir):
            logger.warning(f"Course directory not found: {course_dir}")
            return []
        
        for root, _, files in os.walk(course_dir):
            for file in files:
                file_path = os.path.join(root, file)
                file_ext = os.path.splitext(file)[1].lower()
                
                try:
                    if file_ext == '.pdf':
                        # Extract metadata like titles and page numbers
                        pdf_docs = self._load_pdf_with_metadata(file_path)
                        documents.extend(pdf_docs)
                    elif file_ext == '.pptx':
                        # Extract metadata like slide numbers and titles
                        ppt_docs = self._load_pptx_with_metadata(file_path)
                        documents.extend(ppt_docs)
                    elif file_ext == '.txt':
                        loader = TextLoader(file_path)
                        txt_docs = loader.load()
                        # Add metadata
                        for doc in txt_docs:
                            doc.metadata.update({
                                'source': file_path,
                                'file_name': os.path.basename(file_path),
                                'file_type': 'txt',
                                'title': os.path.basename(file_path),
                                'doc_id': self._generate_doc_id(file_path)
                            })
                        documents.extend(txt_docs)
                    elif file_ext == '.docx':
                        loader = Docx2txtLoader(file_path)
                        docx_docs = loader.load()
                        # Add metadata
                        for doc in docx_docs:
                            doc.metadata.update({
                                'source': file_path,
                                'file_name': os.path.basename(file_path),
                                'file_type': 'docx',
                                'title': os.path.basename(file_path),
                                'doc_id': self._generate_doc_id(file_path)
                            })
                        documents.extend(docx_docs)
                    elif file_ext == '.csv':
                        loader = CSVLoader(file_path)
                        csv_docs = loader.load()
                        # Add metadata
                        for doc in csv_docs:
                            doc.metadata.update({
                                'source': file_path,
                                'file_name': os.path.basename(file_path),
                                'file_type': 'csv',
                                'title': os.path.basename(file_path),
                                'doc_id': self._generate_doc_id(file_path)
                            })
                        documents.extend(csv_docs)
                except Exception as e:
                    logger.error(f"Error loading {file_path}: {str(e)}")
        
        logger.info(f"Loaded {len(documents)} documents for course {self.course_id}")
        return documents

    def _generate_doc_id(self, file_path: str) -> str:
        """Generate a unique ID for a document based on its path"""
        return hashlib.md5(file_path.encode()).hexdigest()

    def _load_pdf_with_metadata(self, file_path: str) -> List[Document]:
        """Load PDF with enhanced metadata including page numbers and titles"""
        documents = []
        try:
            # Use PyMuPDF to extract titles and metadata
            pdf = fitz.open(file_path)
            
            # Get document title from metadata or filename
            doc_title = pdf.metadata.get('title', os.path.basename(file_path))
            
            for page_num in range(len(pdf)):
                page = pdf[page_num]
                text = page.get_text()
                
                # Try to extract page/section title from the first line
                lines = text.strip().split('\n')
                page_title = lines[0].strip() if lines else f"Page {page_num + 1}"
                
                if text.strip():  # Only add non-empty pages
                    doc = Document(
                        page_content=text,
                        metadata={
                            'source': file_path,
                            'file_name': os.path.basename(file_path),
                            'file_type': 'pdf',
                            'page': page_num + 1,
                            'total_pages': len(pdf),
                            'title': doc_title,
                            'page_title': page_title,
                            'doc_id': self._generate_doc_id(file_path)
                        }
                    )
                    documents.append(doc)
            
            pdf.close()
        except Exception as e:
            logger.error(f"Error processing PDF {file_path}: {str(e)}")
        
        return documents

    def _load_pptx_with_metadata(self, file_path: str) -> List[Document]:
        """Load PowerPoint with enhanced metadata including slide numbers and titles"""
        documents = []
        try:
            prs = Presentation(file_path)
            
            for slide_num, slide in enumerate(prs.slides):
                # Extract text from slide
                text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                
                # Try to extract slide title
                slide_title = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.is_title:
                        slide_title = shape.text
                        break
                
                if not slide_title and text.strip():
                    # Use first line as title if no title shape found
                    lines = text.strip().split('\n')
                    slide_title = lines[0] if lines else f"Slide {slide_num + 1}"
                
                if text.strip():  # Only add non-empty slides
                    doc = Document(
                        page_content=text,
                        metadata={
                            'source': file_path,
                            'file_name': os.path.basename(file_path),
                            'file_type': 'pptx',
                            'slide': slide_num + 1,
                            'total_slides': len(prs.slides),
                            'title': os.path.basename(file_path),
                            'slide_title': slide_title,
                            'doc_id': self._generate_doc_id(file_path)
                        }
                    )
                    documents.append(doc)
        except Exception as e:
            logger.error(f"Error processing PowerPoint {file_path}: {str(e)}")
        
        return documents

    def _split_documents(self, documents):
        """Split documents into smaller chunks"""
        logger.info(f"Splitting {len(documents)} documents for course {self.course_id}")
        
        # Create a text splitter with appropriate chunk sizes
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            separators=["\n\n", "\n", ". ", " ", ""],
            keep_separator=False
        )
        
        # Split the documents while preserving metadata
        chunks = text_splitter.split_documents(documents)
        
        # Ensure each chunk has the document metadata
        for i, chunk in enumerate(chunks):
            # Add a chunk ID to help with referencing
            chunk.metadata['chunk_id'] = f"{chunk.metadata.get('doc_id', 'unknown')}_{i}"
        
        logger.info(f"Created {len(chunks)} chunks for course {self.course_id}")
        return chunks
    
    def _initialize_vectorstore(self):
        """Initialize the vector store for the course
        
        This either loads an existing vector store or creates a new one
        from course materials if it doesn't exist.
        """
        # Create vectorstore directory if it doesn't exist
        os.makedirs(self.vectorstore_dir, exist_ok=True)
        
        # Try to load existing vectorstore
        if os.path.exists(self.vectorstore_path):
            logger.info(f"Loading existing vector store for {self.course_id}")
            try:
                with open(self.vectorstore_path, "rb") as f:
                    vectorstore = pickle.load(f)
                # Update embeddings if needed
                vectorstore._embedding = self.embeddings
                return vectorstore
            except Exception as e:
                logger.error(f"Error loading vector store: {str(e)}")
        
        # Create new vectorstore from documents
        logger.info(f"Creating new vector store for {self.course_id}")
        documents = self._load_documents()
        
        if not documents:
            logger.warning("No documents found, creating empty vector store")
            return FAISS.from_documents([], self.embeddings)
        
        chunks = self._split_documents(documents)
        vectorstore = FAISS.from_documents(chunks, self.embeddings)
        
        # Save vectorstore
        logger.info(f"Saving vector store to {self.vectorstore_path}")
        with open(self.vectorstore_path, "wb") as f:
            pickle.dump(vectorstore, f)
        
        return vectorstore
    
    def update_materials(self):
        """Update the vector store with new or modified course materials"""
        logger.info(f"Updating materials for {self.course_id}")
        
        documents = self._load_documents()
        
        if not documents:
            logger.warning("No documents found to update")
            return
        
        chunks = self._split_documents(documents)
        self.vectorstore = FAISS.from_documents(chunks, self.embeddings)
        
        # Save updated vectorstore
        logger.info(f"Saving updated vector store to {self.vectorstore_path}")
        with open(self.vectorstore_path, "wb") as f:
            pickle.dump(self.vectorstore, f)
        
        logger.info("Materials updated successfully")
    
    def _retrieve_context(self, query, top_k=5):
        """Retrieve relevant document chunks for a query"""
        if not hasattr(self, 'vectorstore') or self.vectorstore is None:
            self._initialize_vectorstore()
            
        # Get relevant documents
        docs_with_scores = self.vectorstore.similarity_search_with_score(query, k=top_k)
        
        # Format as context and keep track of references
        context_docs = []
        references = []
        
        for i, (doc, score) in enumerate(docs_with_scores):
            # Create a reference ID
            ref_id = f"ref{i+1}"
            
            # Make sure doc_id exists and is valid
            doc_id = doc.metadata.get('doc_id')
            if not doc_id or doc_id == 'unknown':
                # Generate a doc_id if not present
                if 'source' in doc.metadata:
                    doc_id = self._generate_doc_id(doc.metadata['source'])
                else:
                    # Skip documents with no valid identifiers
                    continue
            
            # Add reference to the document
            ref_info = {
                'id': ref_id,
                'doc_id': doc_id,
                'chunk_id': doc.metadata.get('chunk_id', f"{doc_id}_{i}"),
                'source': doc.metadata.get('source', ''),
                'score': score
            }
            
            # Add page or slide information if available
            if 'page' in doc.metadata:
                ref_info['page'] = doc.metadata['page']
            elif 'slide' in doc.metadata:
                ref_info['slide'] = doc.metadata['slide']
                
            # Add title information
            if 'title' in doc.metadata:
                ref_info['title'] = doc.metadata['title']
            if 'page_title' in doc.metadata:
                ref_info['page_title'] = doc.metadata['page_title']
            elif 'slide_title' in doc.metadata:
                ref_info['slide_title'] = doc.metadata['slide_title']
                
            references.append(ref_info)
            
            # Format document with reference
            context_text = f"{doc.page_content} [Source: {ref_id}]"
            context_docs.append(context_text)
        
        return "\n\n".join(context_docs), references

    def _retrieve_docs_with_scores(self, query: str, k: int = 5):
        if not hasattr(self, 'vectorstore') or self.vectorstore is None:
            self._initialize_vectorstore()
        return self.vectorstore.similarity_search_with_score(query, k=k)
    
    def get_context(self, query, user_id="anonymous", top_k=5):
        """Retrieve context for a query without generating an answer
        
        Args:
            query: The search query
            user_id: Identifier for the user
            top_k: Number of top documents to retrieve
            
        Returns:
            Tuple of (context_text, references)
        """
        logger.info(f"Getting context for query in course {self.course_id}: {query}")
        
        try:
            # Get conversation history
            history = self.conversation_manager.get_conversation_history(self.course_id, user_id)
            
            # Retrieve context documents
            context, references = self._retrieve_context(query, top_k)
            
            return context, references
        except Exception as e:
            logger.error(f"Error getting context: {str(e)}")
            return "", []
    
    def answer_question(self, question: str, user_id: str = "anonymous") -> Tuple[str, List[Dict[str, Any]]]:
        """Agentic answering with a router deciding to use course docs, web, or both.

        Returns (answer, references) with unified refs including web.
        """
        self._ensure_agentic_components()

        logger.info(f"Agentic answering for course {self.course_id}: {question}")

        # History for router
        history = self.conversation_manager.get_conversation_history(self.course_id, user_id)

        # Preview top-3 for router
        docs_with_scores = self._retrieve_docs_with_scores(question, k=3)

        # Route
        decision = self._router.route(
            course_id=self.course_id,
            query=question,
            history=history if isinstance(history, list) else [],
            docs_with_scores=docs_with_scores,
        )

        k_course = max(0, int(decision.get("k_course", 4)))
        k_web = max(0, int(decision.get("k_web", 0)))
        web_queries = decision.get("web_queries", []) if isinstance(decision.get("web_queries", []), list) else []

        # Retrieve course docs (k_course)
        course_docs_with_scores = []
        if k_course > 0:
            course_docs_with_scores = self._retrieve_docs_with_scores(question, k=k_course)
        course_context_texts = []
        course_refs: List[Dict[str, Any]] = []
        for i, (doc, score) in enumerate(course_docs_with_scores):
            ref_id = f"ref{i+1}"
            meta = doc.metadata or {}
            doc_id = meta.get('doc_id') or (self._generate_doc_id(meta['source']) if meta.get('source') else 'unknown')
            ref = {
                'id': ref_id,
                'ref_type': 'course_doc',
                'doc_id': doc_id,
                'chunk_id': meta.get('chunk_id', f"{doc_id}_{i}"),
                'source': meta.get('source', ''),
                'score': score,
            }
            if 'page' in meta:
                ref['page'] = meta['page']
            if 'slide' in meta:
                ref['slide'] = meta['slide']
            if 'title' in meta:
                ref['title'] = meta['title']
            if 'page_title' in meta:
                ref['page_title'] = meta['page_title']
            if 'slide_title' in meta:
                ref['slide_title'] = meta['slide_title']
            course_refs.append(ref)
            course_context_texts.append(f"{doc.page_content} [Source: {ref_id}]")

        web_refs: List[Dict[str, Any]] = []
        web_context_texts: List[str] = []
        if k_web > 0 and web_queries:
            web_results = self._web.search_batch(web_queries, k_each=max(1, min(3, k_web)))
            start_idx = len(course_refs)
            for j, item in enumerate(web_results[:k_web]):
                ref_id = f"ref{start_idx + j + 1}"
                web_ref = {
                    'id': ref_id,
                    'ref_type': 'web',
                    'url': item.get('url', ''),
                    'domain': item.get('domain', ''),
                    'title': item.get('title', ''),
                    'snippet': item.get('snippet', ''),
                    'published_at': item.get('published_at'),
                    'score': item.get('score'),
                    'doc_id': 'web',  # placeholder for UI compatibility
                }
                web_refs.append(web_ref)
                # Keep snippets short to save tokens
                snippet = (item.get('snippet') or '')
                web_context_texts.append(f"{snippet} [Source: {ref_id}]\nURL: {item.get('url','')}")

        # Compose final context (interleave: course first, then web)
        context = "\n\n".join(course_context_texts + web_context_texts)

        # Build prompt and get answer with the existing LLM
        prompt_template = ANSWER_PROMPT

        answer_chain = prompt_template | self.llm | StrOutputParser()
        try:
            answer = answer_chain.invoke(
                {
                    "answer_core_instructions": ANSWER_CORE_INSTRUCTIONS,
                    "course_id": self.course_id,
                    "context": context,
                    "conversation_history": history if history else "No prior conversation.",
                    "query": question,
                }
            )

            references = course_refs + web_refs

            # Persist conversation with refs
            self.conversation_manager.add_message(
                course_id=self.course_id, user_id=user_id, role="user", content=question
            )
            safe_refs = CourseRAG._make_json_serializable(references)
            self.conversation_manager.add_message(
                course_id=self.course_id, user_id=user_id, role="assistant", content=answer, references=safe_refs
            )

            return answer, references
        except Exception as e:
            logger.error(f"Agentic answer error: {str(e)}")
            return "I'm sorry, I ran into an issue answering that. Please try again.", []
        
    def _ensure_agentic_components(self):
        if self._router is None:
            self._router = QueryRouter()
        if self._web is None:
            self._web = WebSearchClient()

    
    def clear_conversation_history(self, user_id):
        """Clear conversation history for a user
        
        Args:
            user_id: Student identifier
            
        Returns:
            True if successful, False otherwise
        """
        return self.conversation_manager.clear_history(self.course_id, user_id)

    def get_document_by_id(self, doc_id: str, page_or_slide: Optional[int] = None) -> Optional[Dict[str, Any]]:
        """Retrieve a document by its ID and optional page/slide number
        
        Args:
            doc_id: The document ID to retrieve
            page_or_slide: The specific page or slide number (optional)
            
        Returns:
            Dictionary with document information or None if not found
        """
        # Ensure vector store is initialized
        if not hasattr(self, 'vectorstore') or self.vectorstore is None:
            self._initialize_vectorstore()
        
        # Get all documents matching the doc_id
        matching_docs = []
        
        # This is a simplified approach; in a production system, you might want 
        # to implement a separate database for more efficient document retrieval
        
        # We need to reload the documents since we don't have direct access to the chunks in the vectorstore
        documents = self._load_documents()
        
        for doc in documents:
            if doc.metadata.get('doc_id') == doc_id:
                if page_or_slide is None or doc.metadata.get('page') == page_or_slide or doc.metadata.get('slide') == page_or_slide:
                    matching_docs.append(doc)
        
        if not matching_docs:
            return None
        
        # Get the first matching document to extract metadata
        doc = matching_docs[0]
        file_path = doc.metadata.get('source')
        file_type = doc.metadata.get('file_type')
        
        # Prepare result
        result = {
            'doc_id': doc_id,
            'file_path': file_path,
            'file_name': os.path.basename(file_path) if file_path else "Unknown",
            'file_type': file_type,
            'title': doc.metadata.get('title', "Untitled Document"),
        }
        
        # Add page/slide specific information
        if file_type == 'pdf':
            result['total_pages'] = doc.metadata.get('total_pages', 0)
            if page_or_slide:
                result['page'] = page_or_slide
        elif file_type == 'pptx':
            result['total_slides'] = doc.metadata.get('total_slides', 0)
            if page_or_slide:
                result['slide'] = page_or_slide
        
        return result


# Example usage
if __name__ == "__main__":
    # Example for testing
    course_id = "DEMO101"
    rag = CourseRAG(course_id)
    
    # Test with a sample question
    question = "What are the key topics covered in this course?"
    print(f"\nQuestion: {question}")
    print(f"\nAnswer: {rag.answer_question(question)}") 