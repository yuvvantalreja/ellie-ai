from flask import Flask, render_template, request, jsonify, send_file, session, Response, redirect, url_for, abort
import os
import threading
import time
import datetime
import uuid
import json
import io
import fitz  # PyMuPDF
from pptx import Presentation
from PIL import Image
import base64
import requests
import numpy as np
from course_rag import CourseRAG
from conversation_manager import ConversationManager
from feedback_system import FeedbackSystem
import sys
# Add the current directory to the path to ensure imports work
sys.path.append(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.abspath(os.path.dirname(__file__)))

try:
    from users import UserManager
except ImportError:
    print("Failed to import UserManager directly, trying alternative import paths")
    # Try different import techniques if the first one fails
    current_dir = os.path.dirname(os.path.abspath(__file__))
    parent_dir = os.path.dirname(current_dir)
    sys.path.insert(0, parent_dir)
    sys.path.insert(0, current_dir)
    from users import UserManager
from firebase_config import initialize_firebase, get_firebase_error_message, login_required, role_required
from dotenv import load_dotenv
from werkzeug.utils import secure_filename

# Load environment variables
load_dotenv()

app = Flask(__name__)
app.secret_key = os.getenv("SECRET_KEY", "ellie-cmu-ai-assistant-secret-key")

# Add template filters
@app.template_filter('now')
def filter_now(format_='%Y'):
    """Return the current year or formatted date"""
    return datetime.datetime.now().strftime(format_)

# Global dict to store RAG instances for different courses
rag_instances = {}
rag_locks = {}  # Locks to prevent concurrent access to RAG instances

conversation_manager = ConversationManager()
feedback_system = FeedbackSystem()
user_manager = UserManager()

# Initialize Firebase
firebase_initialized = initialize_firebase()

# Directory setup
MATERIALS_DIR = os.getenv("MATERIALS_DIR", "course_materials")
UPLOAD_FOLDER = os.path.join("static", "uploads")
os.makedirs(MATERIALS_DIR, exist_ok=True)
os.makedirs("templates", exist_ok=True)
os.makedirs("static", exist_ok=True)
os.makedirs("static/css", exist_ok=True)
os.makedirs("static/js", exist_ok=True)
os.makedirs("static/images", exist_ok=True)
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# Configure appz
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 16MB max upload size

# Get API keys from environment variables
GROQ_API_KEY = os.getenv("GROQ_API_KEY")
GROQ_MODEL = os.getenv("GROQ_MODEL", "llama3-70b-8192")
GROQ_VISION_MODEL = os.getenv("GROQ_VISION_MODEL", "llama-3.2-90b-vision-preview")

def get_rag_instance(course_id, discipline=None):
    """Get or create a RAG instance for the given course"""
    if course_id not in rag_locks:
        rag_locks[course_id] = threading.Lock()
    
    with rag_locks[course_id]:
        if course_id not in rag_instances:
            rag_instances[course_id] = CourseRAG(course_id, discipline=discipline)
        
        return rag_instances[course_id]


def get_user_id():
    """Get the user ID from the session
    
    If the user is authenticated, return their actual user ID.
    Otherwise, create a temporary anonymous ID for the session.
    """
    if 'user_id' in session:
        # Return authenticated user ID
        return session['user_id']


@app.route('/')
def index():
    """Render the main page"""
    # Check if user is logged in
    user_data = None
    is_authenticated = 'user_id' in session
    
    print(f"Is authenticated: {is_authenticated}")
    
    if is_authenticated:
        user_data = user_manager.get_user(session['user_id'])
        courses = []
        
        # Get list of course directories
        if os.path.exists(MATERIALS_DIR):
            courses = [d for d in os.listdir(MATERIALS_DIR) 
                     if os.path.isdir(os.path.join(MATERIALS_DIR, d))]
        
        # If professor, show all courses
        # If student, filter courses they have access to
        if user_data and user_data.get('role') == 'student':
            user_courses = user_data.get('courses', [])
            if user_courses:
                courses = [c for c in courses if c in user_courses]
        
        return render_template('index.html', 
                              courses=courses, 
                              user=user_data, 
                              is_authenticated=is_authenticated)
    else:
        # Redirect unauthenticated users to landing page
        return render_template('landing.html')


@app.route('/landing')
def landing():
    # If user is already logged in, redirect to index
    if 'user_id' in session:
        return redirect(url_for('index'))
    
    return render_template('landing.html')


@app.route('/course/<course_id>')
def course_page(course_id):
    """Render the course-specific chat page"""
    # Ensure user ID is set
    user_id = get_user_id()
    
    # If authenticated user, check if they have access to this course
    if 'user_id' in session:
        user_data = user_manager.get_user(session['user_id'])
        if user_data and user_data.get('role') == 'student':
            user_courses = user_data.get('courses', [])
            if course_id not in user_courses:
                # Add course to user's courses
                user_manager.add_course_to_user(session['user_id'], course_id)
    
    # Get conversation history
    history = conversation_manager.get_conversation_history(course_id, user_id)
    
    # Get user data if authenticated
    user_data = None
    if 'user_id' in session:
        user_data = user_manager.get_user(session['user_id'])
    
    return render_template('chat.html', 
                          course_id=course_id, 
                          history=history,
                          user=user_data,
                          is_authenticated='user_id' in session)


# Authentication Routes
@app.route('/login', methods=['GET', 'POST'])
def login():
    """Handle user login"""
    if 'user_id' in session:
        return redirect(url_for('index'))
    
    error = None
    
    if request.method == 'POST':
        print("Login form submitted")
        email = request.form.get('email')
        password = request.form.get('password')
        
        print(f"Email: {email}")
        print(f"Password length: {len(password) if password else 0}")
        
        if not email or not password:
            error = 'Email and password are required'
        else:
            try:
                from firebase_config import auth
                
                # Check if Firebase is initialized
                if not firebase_initialized:
                    error = 'Authentication service is not available. Please try again later.'
                    print("Firebase not initialized during login attempt")
                    return render_template('login.html', error=error)
                
                print("Attempting to authenticate with Firebase")
                # Authenticate with Firebase
                user = auth.sign_in_with_email_and_password(email, password)
                print(f"Firebase authentication successful: {user['localId']}")
                
                # Get user data from our database
                user_data = user_manager.get_user_by_email(email)
                print(f"User data from database: {user_data}")
                
                if not user_data:
                    error = 'User not found. Please register first.'
                    print("User not found in database")
                else:
                    # Set session data
                    session['user_id'] = user_data['id']
                    session['user_email'] = user_data['email']
                    session['user_name'] = user_data['name']
                    session['user_role'] = user_data['role']
                    
                    print(f"Session data set: {session}")
                    
                    # Update last login
                    user_manager.update_last_login(user_data['id'])
                    
                    # Redirect to home page
                    next_url = request.args.get('next', url_for('index'))
                    print(f"Redirecting to: {next_url}")
                    return redirect(next_url)
                    
            except Exception as e:
                print(f"Login error: {str(e)}")
                error = f'Authentication failed: {str(e)}'
    
    return render_template('login.html', error=error)


@app.route('/register', methods=['GET', 'POST'])
def register():
    """Handle user registration"""
    if 'user_id' in session:
        return redirect(url_for('index'))
    
    error = None
    
    if request.method == 'POST':
        print("Register form submitted")
        name = request.form.get('name')
        email = request.form.get('email')
        password = request.form.get('password')
        role = request.form.get('role')
        
        print(f"Name: {name}")
        print(f"Email: {email}")
        print(f"Password length: {len(password) if password else 0}")
        print(f"Role: {role}")
        
        if not all([name, email, password, role]):
            error = 'All fields are required'
        elif len(password) < 6:
            error = 'Password must be at least 6 characters long'
        else:
            try:
                from firebase_config import auth
                
                # Check if Firebase is initialized
                if not firebase_initialized:
                    error = 'Authentication service is not available. Please try again later.'
                    print("Firebase not initialized during registration attempt")
                    return render_template('register.html', error=error)
                
                print("Attempting to create user in Firebase")
                # Create user in Firebase
                firebase_user = auth.create_user_with_email_and_password(email, password)
                print(f"Firebase user created: {firebase_user['localId']}")
                
                # Create user in our database
                user_data = user_manager.create_user(
                    email=email,
                    name=name,
                    role=role,
                    provider='email',
                    provider_id=firebase_user['localId']
                )
                print(f"User created in database: {user_data}")
                
                # Set session data
                session['user_id'] = user_data['id']
                session['user_email'] = user_data['email']
                session['user_name'] = user_data['name']
                session['user_role'] = user_data['role']
                
                print(f"Session data set: {session}")
                
                # Redirect to home page
                print("Redirecting to home page")
                return redirect(url_for('index'))
                
            except Exception as e:
                error_code = None
                if hasattr(e, 'args') and len(e.args) > 0 and isinstance(e.args[0], dict) and 'error' in e.args[0]:
                    error_code = e.args[0]['error']['message']
                    error = get_firebase_error_message(error_code)
                else:
                    error = f'Registration failed: {str(e)}'
                
                print(f"Registration error: {str(e)}, code: {error_code}")
    
    return render_template('register.html', error=error)


@app.route('/logout')
def logout():
    """Handle user logout"""
    session.pop('user_id', None)
    session.pop('user_email', None)
    session.pop('user_name', None)
    session.pop('user_role', None)
    
    return redirect(url_for('index'))


@app.route('/select_role', methods=['GET', 'POST'])
def select_role():
    """Handle role selection for Google auth users"""
    # This route is only for users who signed in with Google
    if 'temp_user_id' not in session:
        return redirect(url_for('login'))
    
    if request.method == 'POST':
        selected_role = request.form.get('selected_role')
        
        if selected_role not in ['student', 'professor']:
            return render_template('select_role.html', 
                                  name=session.get('temp_user_name'),
                                  error='Please select a valid role')
        
        # Update user role in the database
        user_manager.update_user_role(session['temp_user_id'], selected_role)
        
        # Set permanent session data
        session['user_id'] = session['temp_user_id']
        session['user_email'] = session['temp_user_email']
        session['user_name'] = session['temp_user_name']
        session['user_role'] = selected_role
        
        # Clear temporary data
        session.pop('temp_user_id', None)
        session.pop('temp_user_email', None)
        session.pop('temp_user_name', None)
        
        # Redirect to home page
        return redirect(url_for('index'))
    
    return render_template('select_role.html', name=session.get('temp_user_name'))


# API Routes for Authentication
@app.route('/api/firebase_config', methods=['GET'])
def get_firebase_config():
    """Return Firebase configuration for client-side initialization"""
    from firebase_config import firebase_config
    
    # Only return the essential config that the client needs
    client_config = {
        'apiKey': firebase_config.get('apiKey'),
        'authDomain': firebase_config.get('authDomain'),
        'projectId': firebase_config.get('projectId'),
        'storageBucket': firebase_config.get('storageBucket'),
        'messagingSenderId': firebase_config.get('messagingSenderId'),
        'appId': firebase_config.get('appId')
    }
    
    return jsonify({'success': True, 'config': client_config})


@app.route('/api/auth/google_login', methods=['POST'])
def google_login():
    """Handle Google login API request"""
    print("Google login API endpoint called")
    token = request.json.get('token')
    
    if not token:
        print("Error: Token is required")
        return jsonify({'success': False, 'error': 'Token is required'})
    
    try:
        # Check if Firebase Admin SDK is initialized with service account
        from firebase_config import firebase_admin_app, FIREBASE_SERVICE_ACCOUNT
        if not os.path.exists(FIREBASE_SERVICE_ACCOUNT):
            print(f"Error: Service account file '{FIREBASE_SERVICE_ACCOUNT}' not found")
            return jsonify({
                'success': False, 
                'error': f"Firebase service account file not found. Please create '{FIREBASE_SERVICE_ACCOUNT}' file.",
                'setupRequired': True
            })
        
        if not firebase_admin_app:
            print("Error: Firebase Admin SDK not initialized")
            return jsonify({
                'success': False, 
                'error': 'Firebase Admin SDK not initialized properly',
                'setupRequired': True
            })
        
        # Verify the ID token with Firebase
        from firebase_config import auth
        print(f"Firebase auth object: {auth}")
        
        print("Verifying ID token with Firebase")
        try:
            decoded_token = auth.verify_id_token(token)
            print(f"Decoded token: {decoded_token}")
        except Exception as token_error:
            print(f"Token verification error: {str(token_error)}")
            return jsonify({'success': False, 'error': f'Token verification failed: {str(token_error)}'})
        
        # Get user info
        uid = decoded_token['uid']
        email = decoded_token.get('email', '')
        name = decoded_token.get('name', '')
        if not name:
            name = email.split('@')[0]
        
        print(f"User info from token - UID: {uid}, Email: {email}, Name: {name}")
        
        # Check if user exists in our database
        user_data = user_manager.get_user_by_email(email)
        print(f"User data from database: {user_data}")
        
        if user_data:
            # User exists, set session data
            session['user_id'] = user_data['id']
            session['user_email'] = user_data['email']
            session['user_name'] = user_data['name']
            session['user_role'] = user_data['role']
            
            print(f"Session data set: {session}")
            
            # Update last login
            user_manager.update_last_login(user_data['id'])
            
            print("Redirecting to index page")
            return jsonify({'success': True, 'redirect': url_for('index')})
        else:
            # User doesn't exist, need to complete registration
            print("User doesn't exist, redirecting to complete registration")
            return jsonify({
                'success': True,
                'redirect': url_for('google_complete_registration', token=token)
            })
        
    except Exception as e:
        print(f"Google login error: {str(e)}")
        import traceback
        traceback.print_exc()
        return jsonify({'success': False, 'error': f'Authentication failed: {str(e)}'})


@app.route('/auth/google/complete', methods=['GET'])
def google_complete_registration():
    """Complete registration for Google users"""
    print("Google complete registration route called")
    token = request.args.get('token')
    
    if not token:
        print("No token provided")
        return redirect(url_for('login'))
    
    try:
        # Verify the ID token with Firebase
        from firebase_config import auth
        print("Verifying ID token with Firebase")
        decoded_token = auth.verify_id_token(token)
        print(f"Decoded token: {decoded_token}")
        
        # Get user info
        uid = decoded_token['uid']
        email = decoded_token.get('email', '')
        name = decoded_token.get('name', '')
        if not name:
            name = email.split('@')[0]
        
        print(f"User info - UID: {uid}, Email: {email}, Name: {name}")
        
        # Set temporary session data
        session['temp_user_id'] = str(uuid.uuid4())
        session['temp_user_email'] = email
        session['temp_user_name'] = name
        
        print(f"Temporary session data set: {session}")
        
        # Redirect to role selection
        print("Redirecting to role selection")
        return redirect(url_for('select_role'))
        
    except Exception as e:
        print(f"Google complete registration error: {str(e)}")
        import traceback
        traceback.print_exc()
        return redirect(url_for('login'))


@app.route('/api/auth/google_register', methods=['POST'])
def google_register():
    """Handle Google registration API request"""
    token = request.json.get('token')
    
    if not token:
        return jsonify({'success': False, 'error': 'Token is required'})
    
    try:
        # Check if Firebase Admin SDK is initialized with service account
        from firebase_config import firebase_admin_app, FIREBASE_SERVICE_ACCOUNT
        if not os.path.exists(FIREBASE_SERVICE_ACCOUNT):
            print(f"Error: Service account file '{FIREBASE_SERVICE_ACCOUNT}' not found")
            return jsonify({
                'success': False, 
                'error': f"Firebase service account file not found. Please create '{FIREBASE_SERVICE_ACCOUNT}' file.",
                'setupRequired': True
            })
        
        if not firebase_admin_app:
            print("Error: Firebase Admin SDK not initialized")
            return jsonify({
                'success': False, 
                'error': 'Firebase Admin SDK not initialized properly',
                'setupRequired': True
            })
            
        # Verify the ID token with Firebase
        from firebase_config import auth
        print(f"Google register: Verifying token with Firebase auth")
        
        if auth is None:
            print("Google register error: auth object is None")
            return jsonify({'success': False, 'error': 'Firebase authentication not initialized'})
            
        try:
            decoded_token = auth.verify_id_token(token)
            print(f"Decoded token: {decoded_token}")
        except Exception as token_error:
            print(f"Token verification error: {str(token_error)}")
            return jsonify({'success': False, 'error': f'Token verification failed: {str(token_error)}'})
        
        # Get user info
        uid = decoded_token['uid']
        email = decoded_token.get('email', '')
        name = decoded_token.get('name', '')
        if not name:
            name = email.split('@')[0]
            
        print(f"Google register: User info - UID: {uid}, Email: {email}, Name: {name}")
        
        # Check if user exists in our database
        user_data = user_manager.get_user_by_email(email)
        print(f"Google register: User data from database: {user_data}")
        
        if user_data:
            # User exists, set session data
            print(f"Google register: User exists, setting session data")
            session['user_id'] = user_data['id']
            session['user_email'] = user_data['email']
            session['user_name'] = user_data['name']
            session['user_role'] = user_data['role']
            
            # Update last login
            user_manager.update_last_login(user_data['id'])
            
            return jsonify({'success': True, 'redirect': url_for('index')})
        else:
            # User doesn't exist, need to complete registration
            print(f"Google register: User doesn't exist, redirecting to complete registration")
            return jsonify({
                'success': True,
                'redirect': url_for('google_complete_registration', token=token)
            })
        
    except Exception as e:
        import traceback
        print(f"Google register error: {str(e)}")
        traceback.print_exc()
        return jsonify({'success': False, 'error': f'Authentication failed: {str(e)}'})


@app.route('/api/user/profile', methods=['GET'])
@login_required
def get_user_profile():
    """Get the current user's profile"""
    user_id = session.get('user_id')
    user_data = user_manager.get_user(user_id)
    
    if not user_data:
        return jsonify({'success': False, 'error': 'User not found'})
    
    # Remove sensitive information
    user_data.pop('provider_id', None)
    
    return jsonify({'success': True, 'user': user_data})


@app.route('/api/courses')
def get_courses():
    """API endpoint to get available courses"""
    courses = []
    
    if os.path.exists(MATERIALS_DIR):
        courses = [d for d in os.listdir(MATERIALS_DIR) 
                 if os.path.isdir(os.path.join(MATERIALS_DIR, d))]
    
    # Filter courses for students
    if 'user_id' in session:
        user_data = user_manager.get_user(session['user_id'])
        if user_data and user_data.get('role') == 'student':
            user_courses = user_data.get('courses', [])
            if user_courses:
                courses = [c for c in courses if c in user_courses]
    
    return jsonify({'courses': courses})


# Helper function to make data JSON serializable
def make_json_serializable(obj):
    """Convert NumPy types and other non-serializable types to Python native types"""
    if isinstance(obj, np.float32) or isinstance(obj, np.float64):
        return float(obj)
    elif isinstance(obj, np.int32) or isinstance(obj, np.int64):
        return int(obj)
    elif isinstance(obj, np.ndarray):
        return obj.tolist()
    elif isinstance(obj, dict):
        return {k: make_json_serializable(v) for k, v in obj.items()}
    elif isinstance(obj, list):
        return [make_json_serializable(item) for item in obj]
    else:
        return obj


@app.route('/api/ask', methods=['POST'])
def ask_question():
    """API endpoint to answer a question"""
    data = request.json
    if not data or 'question' not in data or 'course_id' not in data:
        return jsonify({"error": "Missing required parameters"}), 400
    
    question = data['question']
    course_id = data['course_id']
    user_id = get_user_id()
    
    # Get the RAG instance
    rag = get_rag_instance(course_id, data.get('discipline'))
    
    try:
        # Get answer from agentic RAG system (router + web search)
        answer, references = rag.answer_question_agentic(question, user_id)
        
        # Convert references to JSON serializable format
        references = make_json_serializable(references)
        
        # Format the references for the UI
        formatted_refs = []
        for i, ref in enumerate(references):
            ref_text = f"[ref{i+1}]"
            ref_type = ref.get("ref_type")
            if ref_type == "web":
                formatted_refs.append({
                    "id": ref_text,
                    "ref_type": "web",
                    "url": ref.get("url", ""),
                    "domain": ref.get("domain", ""),
                    "title": ref.get("title", ""),
                    "snippet": ref.get("snippet", ""),
                })
            else:
                # Course doc reference
                if not ref.get("doc_id") or ref.get("doc_id") == "unknown":
                    continue
                page_or_slide = None
                if 'page' in ref:
                    page_or_slide = ref['page']
                elif 'slide' in ref:
                    page_or_slide = ref['slide']
                formatted_refs.append({
                    "id": ref_text,
                    "ref_type": "course_doc",
                    "doc_id": ref.get("doc_id", ""),
                    "source": os.path.basename(ref.get("source", "")),
                    "page_or_slide": page_or_slide,
                    "title": ref.get("title", ""),
                    "subtitle": ref.get("page_title", ref.get("slide_title", ""))
                })
        
        # Save the conversation with references
        conversation_manager.add_message(course_id, user_id, "user", question)
        conversation_manager.add_message(course_id, user_id, "assistant", answer, references=references)
        
        return jsonify({
            "answer": answer, 
            "references": formatted_refs
        })
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/api/route_decision_stream', methods=['GET'])
def route_decision_stream():
    """SSE endpoint: stream a quick router decision so UI can reflect web-search intent early."""
    course_id = request.args.get('course_id')
    question = request.args.get('question') or ''
    user_id = get_user_id()
    discipline = request.args.get('discipline')

    if not course_id or not question:
        return jsonify({"error": "Missing required parameters"}), 400

    rag = get_rag_instance(course_id, discipline)

    def generate():
        try:
            # Use internal components to avoid adding new helpers
            rag._ensure_agentic_components()  # type: ignore[attr-defined]
            history = rag.conversation_manager.get_conversation_history(course_id, user_id)
            docs_with_scores = rag._retrieve_docs_with_scores(question, k=3)  # type: ignore[attr-defined]
            decision = rag._router.route(  # type: ignore[attr-defined]
                course_id=course_id,
                query=question,
                history=history if isinstance(history, list) else [],
                docs_with_scores=docs_with_scores,
            )
            k_web = int(decision.get('k_web', 0) or 0)
            used_web = k_web > 0 and bool(decision.get('web_queries') or [])
            payload = json.dumps({"used_web": used_web})
            yield f"data: {payload}\n\n"
            # Indicate completion of the stream
            yield "data: {\"done\": true}\n\n"
        except Exception as e:
            err = json.dumps({"error": str(e)})
            yield f"data: {err}\n\n"

    headers = {
        "Cache-Control": "no-cache",
        "X-Accel-Buffering": "no",  # Disable buffering on some reverse proxies
    }
    return Response(generate(), mimetype='text/event-stream', headers=headers)


@app.route('/api/history', methods=['GET'])
def get_history():
    """API endpoint to get conversation history"""
    course_id = request.args.get('course_id')
    
    if not course_id:
        return jsonify({'error': 'Course ID is required'}), 400
    
    # Get user ID from session
    user_id = get_user_id()
    
    history = conversation_manager.get_conversation_history(course_id, user_id)
    
    return jsonify({'history': history})


@app.route('/api/history/clear', methods=['POST'])
def clear_history():
    """API endpoint to clear conversation history"""
    data = request.json
    course_id = data.get('course_id')
    
    if not course_id:
        return jsonify({'error': 'Course ID is required'}), 400
    
    # Get user ID from session
    user_id = get_user_id()
    
    success = conversation_manager.clear_history(course_id, user_id)
    
    return jsonify({'success': success})



@app.route('/api/update_materials', methods=['POST'])
def update_materials():
    """API endpoint to update course materials"""
    data = request.json
    course_id = data.get('course_id')
    
    if not course_id:
        return jsonify({'error': 'Course ID is required'}), 400
    
    try:
        rag = get_rag_instance(course_id)
        rag.update_materials()
        return jsonify({'success': True, 'message': 'Materials updated successfully'})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/upload_materials', methods=['POST'])
def upload_materials():
    """API endpoint to upload course materials"""
    # Check if course_id is provided
    course_id = request.form.get('course_id')
    if not course_id:
        return jsonify({'error': 'Course ID is required'}), 400
    
    # Check if course directory exists
    course_dir = os.path.join(MATERIALS_DIR, course_id)
    if not os.path.exists(course_dir):
        return jsonify({'error': f'Course {course_id} not found'}), 404
    
    # Check if files were uploaded
    if 'materials' not in request.files:
        return jsonify({'error': 'No files uploaded'}), 400
    
    files = request.files.getlist('materials')
    if not files or files[0].filename == '':
        return jsonify({'error': 'No files selected'}), 400
    
    # Define allowed file extensions
    ALLOWED_EXTENSIONS = {'pdf', 'pptx', 'docx', 'txt', 'csv'}
    
    # Process uploaded files
    uploaded_files = []
    for file in files:
        if file and '.' in file.filename:
            ext = file.filename.rsplit('.', 1)[1].lower()
            if ext in ALLOWED_EXTENSIONS:
                filename = secure_filename(file.filename)
                file_path = os.path.join(course_dir, filename)
                file.save(file_path)
                uploaded_files.append(filename)
    
    if not uploaded_files:
        return jsonify({'error': 'No valid files were uploaded'}), 400
    
    try:
        # Update materials in the RAG system
        rag = get_rag_instance(course_id)
        rag.update_materials()
        
        return jsonify({
            'success': True,
            'message': f'Successfully uploaded {len(uploaded_files)} file(s) to {course_id}',
            'files': uploaded_files
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/create_course', methods=['POST'])
def create_course():
    """API endpoint to create a new course"""
    data = request.json
    course_id = data.get('course_id')
    
    if not course_id:
        return jsonify({'error': 'Course ID is required'}), 400
    
    try:
        # Create course directory
        course_dir = os.path.join(MATERIALS_DIR, course_id)
        os.makedirs(course_dir, exist_ok=True)
        
        return jsonify({
            'success': True, 
            'message': f'Course {course_id} created successfully. Add materials to {course_dir}'
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/document/<course_id>/<doc_id>', methods=['GET'])
def get_document(course_id, doc_id):
    """API endpoint to get document metadata"""
    page_or_slide = request.args.get('page') or request.args.get('slide')
    if page_or_slide:
        try:
            page_or_slide = int(page_or_slide)
        except ValueError:
            page_or_slide = None
    
    rag = get_rag_instance(course_id)
    doc_info = rag.get_document_by_id(doc_id, page_or_slide)
    
    if not doc_info:
        return jsonify({"error": "Document not found"}), 404
    
    return jsonify(doc_info)

@app.route('/api/document/render/<course_id>/<doc_id>', methods=['GET'])
def render_document(course_id, doc_id):
    """Render a document page or slide as an image"""
    # Get page or slide number from query parameters
    page_or_slide = request.args.get('page') or request.args.get('slide')
    if page_or_slide:
        try:
            page_or_slide = int(page_or_slide)
        except ValueError:
            return jsonify({"error": "Invalid page or slide number"}), 400
    
    # Get document info
    rag = get_rag_instance(course_id)
    doc_info = rag.get_document_by_id(doc_id, page_or_slide)
    
    if not doc_info:
        return jsonify({"error": "Document not found"}), 404
    
    file_path = doc_info.get('file_path')
    file_type = doc_info.get('file_type')
    
    if not file_path or not os.path.exists(file_path):
        return jsonify({"error": "Document file not found"}), 404
    
    try:
        # Render based on file type
        if file_type == 'pdf' and page_or_slide:
            # Render PDF page
            pdf = fitz.open(file_path)
            if page_or_slide < 1 or page_or_slide > len(pdf):
                return jsonify({"error": "Page number out of range"}), 400
                
            page = pdf[page_or_slide - 1]  # 0-indexed
            pix = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5))
            img_bytes = pix.tobytes("png")
            
            return Response(img_bytes, mimetype='image/png')
            
        elif file_type == 'pptx' and page_or_slide:
            # Render PowerPoint slide
            prs = Presentation(file_path)
            if page_or_slide < 1 or page_or_slide > len(prs.slides):
                return jsonify({"error": "Slide number out of range"}), 400
            
            # This is a simplification - actual rendering would require
            # converting the slide to an image using additional libraries
            slide = prs.slides[page_or_slide - 1]  # 0-indexed
            
            # Return slide content as JSON for now
            slide_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_content.append(shape.text)
            
            return jsonify({
                "type": "text",
                "content": "\n".join(slide_content)
            })
            
        else:
            # For other file types, just return the file
            return send_file(file_path)
            
    except Exception as e:
        return jsonify({"error": f"Error rendering document: {str(e)}"}), 500

@app.route('/api/document/content/<course_id>/<doc_id>', methods=['GET'])
def get_document_content(course_id, doc_id):
    """Get the text content of a document page or slide"""
    # Get page or slide number from query parameters
    page_or_slide = request.args.get('page') or request.args.get('slide')
    if page_or_slide:
        try:
            page_or_slide = int(page_or_slide)
        except ValueError:
            return jsonify({"error": "Invalid page or slide number"}), 400
    
    # Get document info
    rag = get_rag_instance(course_id)
    doc_info = rag.get_document_by_id(doc_id, page_or_slide)
    
    if not doc_info:
        return jsonify({"error": "Document not found"}), 404
    
    file_path = doc_info.get('file_path')
    file_type = doc_info.get('file_type')
    
    if not file_path or not os.path.exists(file_path):
        return jsonify({"error": "Document file not found"}), 404
    
    try:
        # Extract content based on file type
        if file_type == 'pdf' and page_or_slide:
            # Get PDF page text
            pdf = fitz.open(file_path)
            if page_or_slide < 1 or page_or_slide > len(pdf):
                return jsonify({"error": "Page number out of range"}), 400
                
            page = pdf[page_or_slide - 1]  # 0-indexed
            text = page.get_text()
            
            return jsonify({
                "content": text,
                "title": doc_info.get('title', ''),
                "page": page_or_slide,
                "total_pages": doc_info.get('total_pages', 0)
            })
            
        elif file_type == 'pptx' and page_or_slide:
            # Get PowerPoint slide text
            prs = Presentation(file_path)
            if page_or_slide < 1 or page_or_slide > len(prs.slides):
                return jsonify({"error": "Slide number out of range"}), 400
            
            slide = prs.slides[page_or_slide - 1]  # 0-indexed
            
            # Extract slide title and content
            slide_title = ""
            slide_content = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if shape.is_title:
                        slide_title = shape.text
                    else:
                        slide_content.append(shape.text)
            
            return jsonify({
                "content": "\n".join(slide_content),
                "title": slide_title,
                "slide": page_or_slide,
                "total_slides": doc_info.get('total_slides', 0)
            })
            
        else:
            # For other file types, return an error
            return jsonify({"error": "Content extraction not supported for this file type"}), 400
            
    except Exception as e:
        return jsonify({"error": f"Error extracting content: {str(e)}"}), 500

@app.route('/api/document/download/<course_id>/<doc_id>', methods=['GET'])
def download_document(course_id, doc_id):
    """Download a source document"""
    rag = get_rag_instance(course_id)
    doc_info = rag.get_document_by_id(doc_id)
    
    if not doc_info:
        return jsonify({"error": "Document not found"}), 404
    
    file_path = doc_info.get('file_path')
    
    if not file_path or not os.path.exists(file_path):
        return jsonify({"error": "Document file not found"}), 404
    
    try:
        return send_file(
            file_path,
            as_attachment=True,
            download_name=os.path.basename(file_path)
        )
    except Exception as e:
        return jsonify({"error": f"Error downloading document: {str(e)}"}), 500

@app.route('/api/ask_with_image', methods=['POST'])
def ask_question_with_image():
    """API endpoint to answer a question that may include an image"""
    if not request.form.get('course_id'):
        return jsonify({"error": "Missing course_id parameter"}), 400
    
    course_id = request.form.get('course_id')
    question = request.form.get('question', '')
    user_id = get_user_id()
    
    print(f"Received image upload request for course {course_id}")
    print(f"Files in request: {request.files.keys()}")
    
    # Check if an image was uploaded
    image_file = None
    image_path = None
    if 'image' in request.files and request.files['image'].filename:
        image_file = request.files['image']
        print(f"Image file detected: {image_file.filename}, {image_file.content_type}")
        if image_file and allowed_file(image_file.filename):
            filename = secure_filename(f"{user_id}_{int(time.time())}_{image_file.filename}")
            image_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            image_file.save(image_path)
            print(f"Image saved to: {image_path}")
        else:
            print(f"Invalid image file or not allowed: {image_file.filename}")
    else:
        print("No image file found in request")
    
    try:
        # Get RAG instance for context
        rag = get_rag_instance(course_id)
        
        # If there's an image, use the vision model with the image
        if image_path:
            print(f"Using vision model with image: {image_path}")
            # First, get relevant context from the RAG system based on the text question
            rag_context = ""
            if question.strip():
                rag_context, _ = rag.get_context(question, user_id)
            
            # Prepare prompt with shared core instructions
            try:
                from prompts import ANSWER_CORE_INSTRUCTIONS
                system_prompt = ANSWER_CORE_INSTRUCTIONS.format(course_id=course_id)
            except Exception:
                system_prompt = f"You are Ellie, an AI Teaching Assistant for course {course_id}. Prefer course materials when sufficient; otherwise use available information. Always cite sources as [refN]."
            
            if rag_context:
                system_prompt += f"\n\nHere is relevant information from the course materials that may help you answer:\n{rag_context}"
            
            # Call the vision model with the image
            print(f"Calling vision model {GROQ_VISION_MODEL}")
            answer = call_vision_model(system_prompt, question, image_path)
            print("Vision model response received")
            
            # Save the conversation with the image path
            conversation_manager.add_message(
                course_id, 
                user_id,
                "user", 
                question if question else "[Image uploaded without text]",
                references=[{"image_path": image_path}]
            )
            conversation_manager.add_message(course_id, user_id, "assistant", answer)
            
            return jsonify({
                "answer": answer,
                "references": []  # No direct references for image-based queries yet
            })
        else:
            print("No image path, falling back to text-only response")
            # If no image, use the unified agentic approach
            answer, references = rag.answer_question(question, user_id)
            
            # Convert references to JSON serializable format
            references = make_json_serializable(references)
            
            # Format the references for the UI
            formatted_refs = []
            for i, ref in enumerate(references):
                # Skip references with invalid doc_id
                if not ref.get("doc_id") or ref.get("doc_id") == "unknown":
                    continue
                    
                ref_text = f"[ref{i+1}]"
                page_or_slide = None
                if 'page' in ref:
                    page_or_slide = ref['page']
                elif 'slide' in ref:
                    page_or_slide = ref['slide']
                    
                formatted_refs.append({
                    "id": ref_text,
                    "doc_id": ref.get("doc_id", ""),
                    "source": os.path.basename(ref.get("source", "")),
                    "page_or_slide": page_or_slide,
                    "title": ref.get("title", ""),
                    "subtitle": ref.get("page_title", ref.get("slide_title", ""))
                })
            
            return jsonify({
                "answer": answer,
                "references": formatted_refs
            })
            
    except Exception as e:
        print(f"Error in ask_question_with_image: {str(e)}")
        import traceback
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500

def allowed_file(filename):
    """Check if the uploaded file has an allowed extension"""
    ALLOWED_EXTENSIONS = {'png', 'jpg', 'jpeg', 'gif', 'webp'}
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def call_vision_model(system_prompt, user_question, image_path):
    """Call the GROQ Vision model API with the image and question"""
    if not GROQ_API_KEY:
        raise ValueError("GROQ_API_KEY is not set in environment variables")
    
    print(f"Processing image: {image_path}")
    
    # Get image data as base64
    with open(image_path, "rb") as image_file:
        image_data = base64.b64encode(image_file.read()).decode('utf-8')
    
    print(f"Image encoded as base64 (length: {len(image_data)})")
    
    headers = {
        "Authorization": f"Bearer {GROQ_API_KEY}",
        "Content-Type": "application/json"
    }

    # Combine system prompt with user question since system messages are incompatible with image inputs
    combined_prompt = f"{system_prompt}\n\nQuestion: {user_question if user_question else 'Can you analyze and explain this image?'}"
    
    payload = {
        "model": GROQ_VISION_MODEL,
        "messages": [
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": combined_prompt},
                    {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{image_data}"}}
                ]
            }
        ],
        "temperature": 0.2
    }
    
    print(f"Sending request to Groq Vision API with model: {GROQ_VISION_MODEL}")
    
    response = requests.post(
        "https://api.groq.com/openai/v1/chat/completions",
        headers=headers,
        json=payload
    )
    
    if response.status_code != 200:
        error_msg = f"Error calling GROQ Vision API (status {response.status_code}): {response.text}"
        print(error_msg)
        raise Exception(error_msg)
    
    result = response.json()
    print("Vision API response received successfully")
    
    return result['choices'][0]['message']['content']

if __name__ == '__main__':
    # Create HTML templates and static files
    print("Creating templates directory and HTML files...")
    app.run(debug=True, port=5000) 